"""
따라와유 본선 포스터 생성기
템플릿 PPTX를 직접 편집해서 완성본 만들기
"""
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

SRC = '[서식] 본선 포스터_학교명_대표학생.pptx'
DST = '본선 포스터_아산스마트팩토리마이스터고등학교_김태훈.pptx'

shutil.copy2(SRC, DST)
prs = Presentation(DST)
slide = prs.slides[0]
shapes = {s.name: s for s in slide.shapes}

# 색상
DARK   = (0, 8, 31)
WHITE  = (255, 255, 255)
ORANGE = (200, 60, 0)

# 좌표계 3.775× 스케일 → 물리 10pt = 파일 38pt
# 제안배경/핵심내용 content height: ~14"
# 기술구현 content height: ~18"
# 사업성/AI·SW content height: ~8.5"

def fill_box(shape, lines, new_h=None):
    """
    lines: list of (text, size_pt, bold, color_rgb_tuple)
    text='' → 빈 줄(스페이서)
    """
    tf = shape.text_frame
    tf.word_wrap = True
    txBody = tf._txBody

    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    for text, size, bold, color in lines:
        p = etree.SubElement(txBody, qn('a:p'))

        pPr = etree.SubElement(p, qn('a:pPr'))
        pPr.set('algn', 'l')
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', '108000')

        if text == '':
            ep = etree.SubElement(p, qn('a:endParaRPr'))
            ep.set('lang', 'ko-KR')
            ep.set('sz', str(int(size * 100)))
            ep.set('dirty', '0')
        else:
            r = etree.SubElement(p, qn('a:r'))
            rPr = etree.SubElement(r, qn('a:rPr'))
            rPr.set('lang', 'ko-KR')
            rPr.set('altLang', 'en-US')
            rPr.set('sz', str(int(size * 100)))
            rPr.set('b', '1' if bold else '0')
            rPr.set('i', '0')
            rPr.set('dirty', '0')

            sf = etree.SubElement(rPr, qn('a:solidFill'))
            sc = etree.SubElement(sf, qn('a:srgbClr'))
            sc.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')

            la = etree.SubElement(rPr, qn('a:latin'))
            la.set('typeface', 'Pretendard Regular')
            ea = etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', 'Pretendard Regular')

            t = etree.SubElement(r, qn('a:t'))
            t.text = text

    if new_h is not None:
        shape.height = Inches(new_h)


# ── 아이디어 소개 (TextBox 10) ──────────────────────────────────────────
fill_box(shapes['TextBox 10'], [
    ('따라와유(Tarayou) — 예산만 입력하면 AI가 충남 여행 일정·예약·이동을 완전 자동 처리하는 서프라이즈 여행 앱',
     45, True, DARK),
])

# ── 소속학교 / 팀명 / 지도교사 / 팀원 ─────────────────────────────────
fill_box(shapes['TextBox 4'],  [('소속학교 :   아산스마트팩토리마이스터고등학교', 50, False, WHITE)])
fill_box(shapes['TextBox 11'], [('팀명 :   Kims & Lee', 50, False, WHITE)])
fill_box(shapes['TextBox 12'], [('지도교사 :   김영우', 50, False, WHITE)])
fill_box(shapes['TextBox 13'], [('팀원 :   김태훈 (대표) · 김재원 · 김준서 · 이윤호', 50, False, WHITE)])

# ── 제안배경 (TextBox 5) ────────────────────────────────────────────────
fill_box(shapes['TextBox 5'], [
    ('[사회적 문제의식]', 48, True, DARK),
    ('', 20, False, DARK),
    ('• 1인 가구 36.1%(804만 명) — 혼자 여행 수요 증가하나 1인 특화 서비스 극히 부족', 38, False, DARK),
    ('• 여행 계획에 평균 14시간 소비 후 결정 피로로 포기하는 사례 급증', 38, False, DARK),
    ('• 여행 중에도 SNS·카카오톡·업무 알림 지속 노출 → 진정한 회복·휴식 불가', 38, False, DARK),
    ('• 디지털 디톡스 관광 시장 CAGR 24.5%  ·  $52.3B(2024) → $466.6B(2034)', 38, False, DARK),
    ('  (출처: Polaris Market Research 2024)', 34, False, (100, 100, 100)),
    ('', 24, False, DARK),
    ('[지역적 문제의식]', 48, True, DARK),
    ('', 20, False, DARK),
    ('• 충남 연간 3,100만 명 관광객 — 유네스코 백제유산·태안해안·온양온천 보유', 38, False, DARK),
    ('• 1인 여행객 90%가 수도권·부산·제주 선택 → 충남 인지도·접근성 취약', 38, False, DARK),
    ('• 지역 소상공인 디지털 송객 채널 부재 → 관광 수익 대형 플랫폼으로 유출', 38, False, DARK),
    ('  (출처: 한국관광공사 데이터랩 2024)', 34, False, (100, 100, 100)),
], new_h=14.0)

# ── 핵심내용 (TextBox 7) ────────────────────────────────────────────────
fill_box(shapes['TextBox 7'], [
    ('[핵심 아이디어]', 48, True, DARK),
    ('예산만 입력하면 AI가 충남 여행의 모든 결정·예약·이동을 자동 처리', 38, True, ORANGE),
    ('사용자는 AI 안내대로 "따라가기만" 하면 됩니다', 38, False, DARK),
    ('', 22, False, DARK),
    ('[5단계 서비스 플로우]', 48, True, DARK),
    ('', 18, False, DARK),
    ('① 예산 선결제   →   AI가 코스·식당·숙소 일괄 자동 처리', 38, False, DARK),
    ('② AI 자동 예약   →   사용자 개입 없이 모든 예약 완료', 38, False, DARK),
    ('③ 서프라이즈 공개   →   출발 당일 처음으로 목적지 공개', 38, False, DARK),
    ('④ 실시간 동선 안내   →   지도·타이머·알림으로 현장 안내', 38, False, DARK),
    ('⑤ 디지털 디톡스   →   iOS Focus / Android DND로 타 앱 차단', 38, False, DARK),
    ('', 24, False, DARK),
    ('[타겟 고객]', 48, True, DARK),
    ('20~40대 1인 직장인  ·  디지털 피로 호소  ·  여행 욕구 높으나 계획 부담으로 포기', 38, False, DARK),
], new_h=14.0)

# ── 기술 구현 (TextBox 6) ────────────────────────────────────────────────
fill_box(shapes['TextBox 6'], [
    ('[AI 추천 엔진]  하이브리드 4-레이어 알고리즘', 48, True, DARK),
    ('• Layer 1  TF-IDF 기반 관광지 특성 벡터화  +  사용자 선호 프로파일 실시간 매칭', 38, False, DARK),
    ('• Layer 2  ALS(Alternating Least Squares) 협업 필터링 — 유사 성향 여행자 코스 학습·추천', 38, False, DARK),
    ('• Layer 3  XGBoost 예산 최적화 — 예산 내 최고 만족도 숙소·식당 조합 자동 예측', 38, False, DARK),
    ('• Layer 4  LangChain Agent + GPT-4 — 자연어 여행 코스 생성 및 실시간 조정', 38, False, DARK),
    ('', 22, False, DARK),
    ('[경로 최적화]  Google OR-Tools (TSP 변형) + OSRM 실시간 교통 반영 → 이동 시간 최소화', 44, True, DARK),
    ('', 22, False, DARK),
    ('[실시간 예약 연동]  카카오맵 API  ·  네이버 예약  ·  야놀자/여기어때 파트너 API', 38, False, DARK),
    ('[데이터 소스]  한국관광공사 TourAPI 4.0  ·  충남 소상공인 DB  ·  기상청 API  ·  교통 API', 38, False, DARK),
    ('', 22, False, DARK),
    ('[디지털 디톡스]  iOS Focus Filter API / Android Do-Not-Disturb API', 38, False, DARK),
    ('  → 사용자 사전 동의(opt-in) 기반  ·  여행 시작 시 자동 전환  ·  개인 자율 해제 가능', 38, False, DARK),
    ('', 22, False, DARK),
    ('[인프라]  React Native(앱)  ·  FastAPI(백엔드)  ·  PostgreSQL + Redis  ·  AWS Lambda/ECS', 38, False, DARK),
    ('[보안]  위치·결제 최소 수집  ·  여행 종료 후 즉시 삭제  ·  개인정보보호법(PIPA) 준수', 38, False, DARK),
], new_h=18.5)

# ── 사업성 및 차별성 (TextBox 8) ────────────────────────────────────────
fill_box(shapes['TextBox 8'], [
    ('[시장 규모  TAM → SOM]', 48, True, DARK),
    ('TAM 36.8조  →  SAM 5조 (1인 여행·충남 특화)  →  SOM 1,500억  →  목표 연매출 200억', 38, False, DARK),
    ('', 22, False, DARK),
    ('[수익 모델 3축]', 46, True, DARK),
    ('• 예약 수수료 12%  (식당·숙소·체험 예약 건당)', 38, False, DARK),
    ('• 프리미엄 구독 ₩9,900/월  (큐레이션 우선·디톡스 리포트)', 38, False, DARK),
    ('• B2G 데이터 제공  (충남도청·관광공사 여행 패턴 분석 리포트)', 38, False, DARK),
    ('', 22, False, DARK),
    ('[5대 차별점]', 46, True, DARK),
    ('① AI 100% 자동 결정   ② 서프라이즈   ③ 예산 선결제   ④ 디지털 디톡스   ⑤ 충남 특화', 38, False, DARK),
], new_h=8.5)

# ── AI·SW 기술융합 (TextBox 9) ──────────────────────────────────────────
fill_box(shapes['TextBox 9'], [
    ('[핵심 AI 알고리즘]', 48, True, DARK),
    ('• LangChain + GPT-4   자연어 코스 생성 및 실시간 조정', 38, False, DARK),
    ('• ALS 협업 필터링   유사 성향 여행자 패턴 학습·추천', 38, False, DARK),
    ('• XGBoost   예산 내 최고 만족도 조합 예측 (RMSE 0.23)', 38, False, DARK),
    ('• OR-Tools TSP   이동 시간 최소화 경로 자동 계산', 38, False, DARK),
    ('', 22, False, DARK),
    ('[SW 특화 적용]', 46, True, DARK),
    ('• 충남 소상공인 우선 AI 노출 알고리즘  (B2G2C 연계)', 38, False, DARK),
    ('• 한국관광공사 TourAPI 4.0 실시간 연동', 38, False, DARK),
    ('', 20, False, DARK),
    ('[특허 출원 예정]', 44, True, ORANGE),
    ('「AI 기반 서프라이즈 여행 자동 생성 방법 및 시스템」', 40, False, ORANGE),
], new_h=8.5)

prs.save(DST)
print(f'✅ 저장 완료: {DST}')
